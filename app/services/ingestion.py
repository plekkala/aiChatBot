"""
Document Ingestion Service
Handles: file parsing → text extraction → chunking → embedding → DB storage
"""
import io
import os
from pathlib import Path
from typing import List, Tuple

import markdown
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from sqlalchemy.orm import Session

from app.core.config import settings
from app.db.models import Document, DocumentChunk

# Load embedding model once at import time
_embedder = SentenceTransformer("all-MiniLM-L6-v2")

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


# ─── Text Extraction ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> List[Tuple[int, str]]:
    """Returns list of (page_number, text) tuples."""
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((i + 1, text.strip()))
    return pages


def extract_text_from_docx(file_bytes: bytes) -> List[Tuple[int, str]]:
    doc = DocxDocument(io.BytesIO(file_bytes))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [(1, full_text)]


def extract_text_from_pptx(file_bytes: bytes) -> List[Tuple[int, str]]:
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append((i + 1, "\n".join(texts)))
    return slides


def extract_text_from_txt(file_bytes: bytes) -> List[Tuple[int, str]]:
    text = file_bytes.decode("utf-8", errors="replace")
    return [(1, text.strip())]


def extract_text_from_md(file_bytes: bytes) -> List[Tuple[int, str]]:
    raw = file_bytes.decode("utf-8", errors="replace")
    # Strip HTML tags after markdown conversion for clean text
    import re
    html = markdown.markdown(raw)
    clean = re.sub(r"<[^>]+>", " ", html).strip()
    return [(1, clean)]


def extract_text(filename: str, file_bytes: bytes) -> List[Tuple[int, str]]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == ".md":
        return extract_text_from_md(file_bytes)
    elif ext == ".txt":
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ─── Chunking ─────────────────────────────────────────────────────────────────

def chunk_text(text: str) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)


# ─── Embedding ────────────────────────────────────────────────────────────────

def embed_texts(texts: List[str]) -> List[List[float]]:
    vectors = _embedder.encode(texts, show_progress_bar=False, normalize_embeddings=True)
    return vectors.tolist()


# ─── Main Ingestion Entry Point ───────────────────────────────────────────────

def ingest_document(db: Session, filename: str, file_bytes: bytes) -> Document:
    """
    Full pipeline: extract → chunk → embed → persist.
    Returns the saved Document record.
    """
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file extension '{ext}'. Supported: {SUPPORTED_EXTENSIONS}")

    # 1. Extract text (list of page/slide tuples)
    pages = extract_text(filename, file_bytes)
    if not pages:
        raise ValueError("No text content could be extracted from the document.")

    # 2. Persist Document record
    doc = Document(
        filename=filename,
        file_type=ext.lstrip("."),
        title=Path(filename).stem.replace("_", " ").replace("-", " ").title(),
    )
    db.add(doc)
    db.flush()  # get doc.id before bulk insert

    # 3. Chunk, embed, and store each page
    chunk_records = []
    global_index = 0

    for page_num, page_text in pages:
        chunks = chunk_text(page_text)
        if not chunks:
            continue
        vectors = embed_texts(chunks)
        for chunk_text_content, vector in zip(chunks, vectors):
            chunk_records.append(
                DocumentChunk(
                    document_id=doc.id,
                    chunk_index=global_index,
                    content=chunk_text_content,
                    embedding=vector,
                    page_number=page_num,
                )
            )
            global_index += 1

    db.bulk_save_objects(chunk_records)
    db.commit()
    db.refresh(doc)

    print(f"[Ingest] '{filename}' → {global_index} chunks stored (doc_id={doc.id})")
    return doc
